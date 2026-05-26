"""
Pitch Deck PowerPoint Generator Service

Generates professional PPTX files from slide data using python-pptx.
Supports dark / light / blue themes. Uses real native python-pptx charts.

Improvements applied:
  1. Real native PIE + BAR charts (XL_CHART_TYPE) — no fake shapes
  2. Dynamic text box height — bullets never truncate
  3. Auto logo fetch — multi-source internet search by company name:
       Layer 1: Clearbit Logo API          (logo.clearbit.com)
       Layer 2: Google Favicon HD          (www.google.com/s2/favicons?sz=256)
       Layer 3: Brandfetch Community API   (api.brandfetch.io/v2/search)
       Layer 4: DuckDuckGo favicon         (icons.duckduckgo.com/ip3)
       Falls back gracefully — never crashes deck build.

Requires:
    pip install python-pptx pillow requests lxml
"""

import io
import re
import requests
from PIL import Image as PILImage

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.oxml.ns import qn
from lxml import etree

# ── Slide dimensions (16:9 widescreen) ───────────────────────────────────────
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

PT_PER_INCH  = 72
EMU_PER_INCH = 914400


# ── Theme definitions ─────────────────────────────────────────────────────────
THEMES = {
    "dark": {
        "bg":           RGBColor(13,  17,  23),
        "card":         RGBColor(22,  27,  34),
        "accent":       RGBColor(35,  134, 54),
        "accent2":      RGBColor(31,  136, 161),
        "text":         RGBColor(240, 246, 252),
        "muted":        RGBColor(110, 118, 129),
        "header_text":  RGBColor(35,  134, 54),
        "chart_colors": ["238636", "1F88A1", "6E7681", "FFB400", "DC5050"],
    },
    "light": {
        "bg":           RGBColor(255, 255, 255),
        "card":         RGBColor(241, 245, 249),
        "accent":       RGBColor(37,  99,  235),
        "accent2":      RGBColor(6,   182, 212),
        "text":         RGBColor(15,  23,  42),
        "muted":        RGBColor(100, 116, 139),
        "header_text":  RGBColor(255, 255, 255),
        "chart_colors": ["2563EB", "06B6D4", "64748B", "F59E0B", "EF4444"],
    },
    "blue": {
        "bg":           RGBColor(15,  35,  68),
        "card":         RGBColor(23,  54,  105),
        "accent":       RGBColor(59,  130, 246),
        "accent2":      RGBColor(6,   182, 212),
        "text":         RGBColor(226, 232, 240),
        "muted":        RGBColor(148, 163, 184),
        "header_text":  RGBColor(255, 255, 255),
        "chart_colors": ["3B82F6", "06B6D4", "94A3B8", "F59E0B", "EF4444"],
    },
}


def _get_theme(context: dict) -> dict:
    key = (context or {}).get("theme", "dark").lower()
    return THEMES.get(key, THEMES["dark"])


# ── IMPROVEMENT 3 — Multi-source Auto Logo Fetch ─────────────────────────────

def _name_to_slug(startup_name: str) -> tuple:
    """
    Convert 'Startup Name Inc.' → slug='startupname', readable='startup-name'
    Returns (slug, readable_slug) for domain guessing.
    """
    clean = startup_name.lower().strip()
    for suffix in [' inc', ' ltd', ' llc', ' corp', ' technologies',
                   ' technology', ' solutions', ' services', ',', '.']:
        clean = clean.replace(suffix, '')
    clean = clean.strip()
    slug     = re.sub(r'[^a-z0-9]', '', clean)
    readable = re.sub(r'[^a-z0-9]+', '-', clean).strip('-')
    return slug, readable


def _is_valid_image(content: bytes, min_size: int = 32) -> bool:
    """Return True if content is a decodable image >= min_size x min_size px."""
    try:
        img = PILImage.open(io.BytesIO(content))
        w, h = img.size
        return w >= min_size and h >= min_size
    except Exception:
        return False


def _try_url(url: str, timeout: int = 5) -> bytes | None:
    """
    Fetch URL, return raw bytes if it's a valid image, else None.
    Never raises — all exceptions swallowed silently.
    """
    try:
        headers = {
            "User-Agent": (
                "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
                "AppleWebKit/537.36 (KHTML, like Gecko) "
                "Chrome/124.0 Safari/537.36"
            )
        }
        r = requests.get(url, timeout=timeout, headers=headers)
        if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
            if _is_valid_image(r.content):
                return r.content
    except Exception:
        pass
    return None


def _fetch_logo(startup_name: str) -> bytes | None:
    """
    Fetch company logo by searching the internet via 4 layered sources.

    Layer 1 — Clearbit Logo API (best quality, works for top 500K companies)
        https://logo.clearbit.com/{domain}
        Free, no key needed, returns 200x200 PNG.

    Layer 2 — Google S2 Favicon Service (high-res 256px favicons)
        https://www.google.com/s2/favicons?sz=256&domain={domain}
        Works for almost every website that has a favicon set.
        Often returns the actual logo PNG at 256px.

    Layer 3 — Brandfetch Community API (brand logo database)
        https://api.brandfetch.io/v2/search?c=1&query={name}
        Returns JSON with logo CDN URLs — parses first PNG/SVG result.
        Free community tier, no API key required for search.

    Layer 4 — DuckDuckGo favicon proxy (universal fallback)
        https://icons.duckduckgo.com/ip3/{domain}.ico
        DDG caches favicons for billions of domains — very high hit rate.
        Lower quality (ico/32px) but extremely reliable.

    Returns PNG/JPEG bytes or None — never crashes the deck build.
    """
    if not startup_name:
        return None

    slug, readable = _name_to_slug(startup_name)

    # Domain candidates to try across all layers
    domain_candidates = [
        f"{slug}.com",
        f"{slug}.io",
        f"{slug}.ai",
        f"{slug}.in",
        f"{readable}.com",
        f"{readable}.io",
    ]

    # ── Layer 1: Clearbit Logo API ────────────────────────────────────────────
    for domain in domain_candidates:
        data = _try_url(f"https://logo.clearbit.com/{domain}")
        if data:
            print(f"[pptx] ✓ Logo via Clearbit: {domain}")
            return data

    # ── Layer 2: Google S2 Favicon (256px) ────────────────────────────────────
    # Google's favicon service often returns crisp logo PNGs at 256px
    for domain in domain_candidates:
        url  = f"https://www.google.com/s2/favicons?sz=256&domain={domain}"
        data = _try_url(url)
        if data and _is_valid_image(data, min_size=64):
            # Skip 16x16 default "globe" favicon Google returns for unknown domains
            img  = PILImage.open(io.BytesIO(data))
            w, h = img.size
            if w >= 64 and h >= 64:
                print(f"[pptx] ✓ Logo via Google Favicon: {domain} ({w}x{h})")
                return data

    # ── Layer 3: Brandfetch Community Search API ───────────────────────────────
    # Returns JSON list of brand objects with logo CDN URLs
    try:
        bf_url = f"https://api.brandfetch.io/v2/search?c=1&query={requests.utils.quote(startup_name)}"
        r = requests.get(bf_url, timeout=6, headers={
            "User-Agent": "Mozilla/5.0",
            "Accept": "application/json",
        })
        if r.status_code == 200:
            results = r.json()
            if isinstance(results, list) and results:
                for brand in results[:3]:
                    # Each brand may have an 'icon' or 'logo' field with CDN URL
                    icon_url = (brand.get("icon") or brand.get("logo") or "").strip()
                    if icon_url and icon_url.startswith("http"):
                        data = _try_url(icon_url)
                        if data:
                            print(f"[pptx] ✓ Logo via Brandfetch: {brand.get('name', startup_name)}")
                            return data
                    # Also check nested formats: brand.logos[0].formats[0].src
                    for logo_obj in brand.get("logos", [])[:2]:
                        for fmt in logo_obj.get("formats", [])[:3]:
                            src = fmt.get("src", "")
                            if src and src.startswith("http"):
                                data = _try_url(src)
                                if data:
                                    print(f"[pptx] ✓ Logo via Brandfetch nested: {startup_name}")
                                    return data
    except Exception as e:
        print(f"[pptx] Brandfetch layer skipped: {e}")

    # ── Layer 4: DuckDuckGo favicon proxy (universal fallback) ────────────────
    # DDG caches favicons for billions of sites — very high hit rate
    for domain in domain_candidates:
        data = _try_url(f"https://icons.duckduckgo.com/ip3/{domain}.ico")
        if data and _is_valid_image(data, min_size=32):
            img  = PILImage.open(io.BytesIO(data))
            w, h = img.size
            # Upscale small favicons with PIL so they don't look blurry in PPTX
            if w < 128:
                scale = 128 // w
                img   = img.resize((w * scale, h * scale), PILImage.LANCZOS)
                buf   = io.BytesIO()
                img.save(buf, format="PNG")
                data  = buf.getvalue()
            print(f"[pptx] ✓ Logo via DuckDuckGo favicon: {domain}")
            return data

    print(f"[pptx] ⚠  No logo found for '{startup_name}' across all 4 sources — skipping")
    return None


# ── Low-level shape helpers ───────────────────────────────────────────────────

def _blank(prs: Presentation):
    """Always use blank layout (index 6) — avoids placeholder KeyError crashes."""
    return prs.slides.add_slide(prs.slide_layouts[6])


def _bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _rect(slide, left, top, width, height,
          fill_color: RGBColor, line_color=None):
    shp = slide.shapes.add_shape(1, left, top, width, height)
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill_color
    if line_color:
        shp.line.color.rgb = line_color
    else:
        shp.line.fill.background()
    return shp


def _textbox(slide, left, top, width, height,
             text: str, size: int, bold=False,
             color: RGBColor = None,
             align=PP_ALIGN.LEFT) -> None:
    txb = slide.shapes.add_textbox(left, top, width, height)
    tf  = txb.text_frame
    tf.word_wrap = True
    lines = str(text).split('\n')
    for li, line in enumerate(lines):
        p = tf.paragraphs[0] if li == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run()
        run.text           = line.strip()
        run.font.size      = Pt(size)
        run.font.bold      = bold
        run.font.color.rgb = color or RGBColor(240, 246, 252)
        run.font.name      = "Calibri"


# ── IMPROVEMENT 2 — Dynamic text height ──────────────────────────────────────

def _dynamic_bullet_height(text: str, font_size_pt: int,
                            box_width_inches: float,
                            line_padding_inches: float = 0.12) -> float:
    """
    Estimate rendered height of bullet text to prevent truncation.
    Counts wrapped lines based on average chars per inch for Calibri.
    Returns height in inches (minimum 0.55").
    """
    chars_per_inch = 9.5 * (12 / font_size_pt)
    chars_per_line = max(int(box_width_inches * chars_per_inch), 1)

    words   = text.split()
    lines   = 0
    cur_len = 0
    for word in words:
        if cur_len + len(word) + 1 > chars_per_line:
            lines  += 1
            cur_len = len(word)
        else:
            cur_len += len(word) + 1
    lines += 1  # last line

    font_h = font_size_pt / PT_PER_INCH
    return max(lines * (font_h + line_padding_inches), 0.55)


# ── IMPROVEMENT 1 — Real native python-pptx charts ───────────────────────────

def _add_real_chart(slide, chart_data: dict, theme: dict,
                    left, top, width, height):
    """
    Build a REAL native Office chart object from LLM chart_data.

    chart_data keys (from pitch_ai.py LLM output):
        type   : "pie" | "bar"
        title  : chart title string
        labels : list of category labels   e.g. ["TAM", "SAM", "SOM"]
        values : list of numeric values    e.g. [120, 18, 1.8]

    For 'pie' → XL_CHART_TYPE.PIE  with coloured slices + % labels
    For 'bar' → XL_CHART_TYPE.BAR_CLUSTERED with coloured bars + value labels
    """
    labels = chart_data.get("labels", [])
    values = chart_data.get("values", [])
    title  = chart_data.get("title",  "")
    ctype  = chart_data.get("type",   "pie")

    if not labels or not values:
        return

    # Sanitise values
    safe_values = []
    for v in values:
        try:
            safe_values.append(float(v))
        except (TypeError, ValueError):
            safe_values.append(0.0)

    n = min(len(labels), len(safe_values))
    labels      = labels[:n]
    safe_values = safe_values[:n]

    cd = ChartData()
    cd.categories = labels
    cd.add_series("", safe_values)

    xl_type = XL_CHART_TYPE.PIE if ctype == "pie" else XL_CHART_TYPE.BAR_CLUSTERED

    chart_shape = slide.shapes.add_chart(xl_type, left, top, width, height, cd)
    chart = chart_shape.chart

    # ── Chart title ───────────────────────────────────────────────────────────
    chart.has_title = True
    chart.chart_title.text_frame.text = title
    for para in chart.chart_title.text_frame.paragraphs:
        for run in para.runs:
            run.font.color.rgb = theme["accent"]
            run.font.bold      = True
            run.font.size      = Pt(13)
            run.font.name      = "Calibri"

    # ── Legend ────────────────────────────────────────────────────────────────
    chart.has_legend = True
    chart.legend.font.size      = Pt(10)
    chart.legend.font.color.rgb = theme["muted"]
    chart.legend.font.name      = "Calibri"

    # ── Per-slice / per-bar colours via XML ───────────────────────────────────
    colors = theme["chart_colors"]
    try:
        ser = chart._element.find('.//' + qn('c:ser'))
        if ser is not None:
            for idx in range(len(labels)):
                hex_val = colors[idx % len(colors)]
                dpt = etree.fromstring(
                    f'<c:dPt xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
                    f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                    f'<c:idx val="{idx}"/><c:invertIfNegative val="0"/>'
                    f'<c:spPr><a:solidFill><a:srgbClr val="{hex_val}"/></a:solidFill>'
                    f'<a:ln><a:noFill/></a:ln></c:spPr></c:dPt>'
                )
                insert_pos = 2
                ser.insert(insert_pos + idx, dpt)
    except Exception as e:
        print(f"[pptx] Chart colour warning: {e}")

    # ── Data labels ───────────────────────────────────────────────────────────
    try:
        for ser_el in chart._element.findall('.//' + qn('c:ser')):
            label_type = 'showPercent' if ctype == 'pie' else 'showVal'
            dLbls = etree.fromstring(
                f'<c:dLbls xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
                f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<c:txPr><a:bodyPr/><a:lstStyle/>'
                f'<a:p><a:pPr><a:defRPr sz="1000" b="1">'
                f'<a:solidFill><a:srgbClr val="F0F6FC"/></a:solidFill>'
                f'</a:defRPr></a:pPr></a:p></c:txPr>'
                f'<c:showLegendKey val="0"/>'
                f'<c:showVal val="{"1" if ctype != "pie" else "0"}"/>'
                f'<c:showCatName val="{"1" if ctype == "pie" else "0"}"/>'
                f'<c:showSerName val="0"/>'
                f'<c:showPercent val="{"1" if ctype == "pie" else "0"}"/>'
                f'<c:showBubbleSize val="0"/>'
                f'</c:dLbls>'
            )
            existing = ser_el.find(qn('c:dLbls'))
            if existing is not None:
                ser_el.remove(existing)
            ser_el.append(dLbls)
    except Exception as e:
        print(f"[pptx] Data label warning: {e}")

    # ── Dark background for chart area ────────────────────────────────────────
    try:
        bg_hex = f"{theme['card'].red:02X}{theme['card'].green:02X}{theme['card'].blue:02X}"
        plot_area = chart._element.find('.//' + qn('c:plotArea'))
        if plot_area is not None:
            pa_spPr = etree.fromstring(
                f'<c:spPr xmlns:c="http://schemas.openxmlformats.org/drawingml/2006/chart"'
                f' xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
                f'<a:solidFill><a:srgbClr val="{bg_hex}"/></a:solidFill>'
                f'<a:ln><a:noFill/></a:ln></c:spPr>'
            )
            existing_sp = plot_area.find(qn('c:spPr'))
            if existing_sp is not None:
                plot_area.remove(existing_sp)
            plot_area.append(pa_spPr)
    except Exception as e:
        print(f"[pptx] Chart bg warning: {e}")


# ── Slide builders ────────────────────────────────────────────────────────────

def _build_title_slide(prs: Presentation, slide_data: dict, theme: dict,
                       startup_name: str, logo_bytes):
    """Cover slide — dark/light/blue bg + company name + tagline + auto logo."""
    slide = _blank(prs)
    _bg(slide, theme["bg"])

    # Bottom strip
    _rect(slide, Inches(0), Inches(6.0), SLIDE_W, Inches(1.5), theme["card"])
    # Left accent bar
    _rect(slide, Inches(0), Inches(0), Inches(0.18), SLIDE_H, theme["accent"])
    # Horizontal divider line
    _rect(slide, Inches(0), Inches(5.95), SLIDE_W, Inches(0.07), theme["accent2"])

    # ── Logo — top-right corner, beside company name ──────────────────────────
    if logo_bytes:
        try:
            img    = PILImage.open(io.BytesIO(logo_bytes))
            # Convert to RGBA to handle transparent PNGs cleanly
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGBA")
            iw, ih = img.size
            # Target: max 1.8" wide, 0.9" tall — keeps logo proportional
            max_w  = Inches(1.8)
            max_h  = Inches(0.9)
            scale  = min(max_w / iw, max_h / ih)
            fw     = int(iw * scale)
            fh     = int(ih * scale)
            # Re-encode as PNG to preserve transparency
            buf = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            # Place at top-right: right-aligned with 0.25" margin, vertically centred in 1.4" top strip
            left_pos = int(SLIDE_W - fw - Inches(0.4))
            top_pos  = int((Inches(1.4) - fh) / 2)
            slide.shapes.add_picture(buf, left=left_pos, top=top_pos, width=fw, height=fh)
            print(f"[pptx] ✓ Logo embedded on cover ({iw}x{ih} → scaled)")
        except Exception as e:
            print(f"[pptx] Cover logo skipped: {e}")

    title_text   = slide_data.get("title", startup_name)
    content_text = slide_data.get("content", "")
    if isinstance(content_text, list):
        content_text = " ".join(str(i) for i in content_text)

    # Company name — shifted left to leave room for logo on the right
    _textbox(slide,
             left=Inches(0.55), top=Inches(1.6),
             width=Inches(10.5), height=Inches(1.9),
             text=title_text, size=54, bold=True,
             color=theme["accent"], align=PP_ALIGN.LEFT)

    # Tagline
    _textbox(slide,
             left=Inches(0.55), top=Inches(3.6),
             width=Inches(11.0), height=Inches(1.5),
             text=content_text, size=24,
             color=theme["text"], align=PP_ALIGN.LEFT)

    # Footer
    _textbox(slide,
             left=Inches(0.55), top=Inches(6.2),
             width=Inches(9), height=Inches(0.5),
             text="CONFIDENTIAL  ·  INVESTOR PRESENTATION",
             size=10, color=theme["muted"], align=PP_ALIGN.LEFT)

    _textbox(slide,
             left=Inches(10.5), top=Inches(6.2),
             width=Inches(2.5), height=Inches(0.5),
             text="2026", size=10,
             color=theme["muted"], align=PP_ALIGN.RIGHT)


def _build_content_slide(prs: Presentation, slide_data: dict,
                         slide_num: int, theme: dict, startup_name: str):
    """
    Standard content slide.
    - Bullets on left (full width if no chart, 55% if chart present)
    - REAL native chart on right 43% when chart_data present
    - Dynamic bullet height — content never truncates
    """
    slide = _blank(prs)
    _bg(slide, theme["bg"])

    title_text   = slide_data.get("title", f"Slide {slide_num}")
    content_text = slide_data.get("content", "")
    if isinstance(content_text, list):
        content_text = "\n".join(str(i) for i in content_text)
    chart_data = slide_data.get("chart_data")
    has_chart  = bool(chart_data and chart_data.get("values"))

    # ── Header band ───────────────────────────────────────────────────────────
    _rect(slide, Inches(0),    Inches(0), SLIDE_W,      Inches(1.1), theme["card"])
    _rect(slide, Inches(0),    Inches(0), Inches(0.18), Inches(1.1), theme["accent"])

    _textbox(slide,
             left=Inches(0.35), top=Inches(0.1),
             width=Inches(11.8), height=Inches(0.9),
             text=title_text, size=30, bold=True,
             color=theme["accent"], align=PP_ALIGN.LEFT)

    _textbox(slide,
             left=Inches(12.4), top=Inches(0.15),
             width=Inches(0.8), height=Inches(0.5),
             text=str(slide_num), size=10,
             color=theme["muted"], align=PP_ALIGN.RIGHT)

    # ── Bullet content with dynamic height ────────────────────────────────────
    bullet_right_edge = 6.8 if has_chart else 12.8
    bullet_width_in   = bullet_right_edge - 0.85

    lines   = [l.strip() for l in content_text.split("\n") if l.strip()]
    top     = Inches(1.25)
    font_sz = 17
    padding = Inches(0.18)

    for line in lines[:7]:
        dyn_h  = _dynamic_bullet_height(line, font_sz, bullet_width_in)
        box_h  = Inches(dyn_h)
        dot_top = top + Inches((font_sz / PT_PER_INCH) / 2 + 0.04)

        # Bullet dot
        _rect(slide,
              left=Inches(0.32), top=dot_top,
              width=Inches(0.14), height=Inches(0.14),
              fill_color=theme["accent"])

        # Bullet text
        _textbox(slide,
                 left=Inches(0.62), top=top,
                 width=Inches(bullet_width_in), height=box_h,
                 text=line, size=font_sz, color=theme["text"])

        top += box_h + padding

    # ── Footer watermark ──────────────────────────────────────────────────────
    _textbox(slide,
             left=Inches(0.35), top=Inches(7.1),
             width=Inches(6), height=Inches(0.3),
             text=startup_name, size=9,
             color=theme["muted"], align=PP_ALIGN.LEFT)

    # ── REAL native chart (right side) ────────────────────────────────────────
    if has_chart:
        _add_real_chart(
            slide, chart_data, theme,
            left=Inches(7.1), top=Inches(1.15),
            width=Inches(5.9), height=Inches(5.9)
        )


def _build_closing_slide(prs: Presentation, theme: dict,
                         startup_name: str, logo_bytes):
    """Thank You / closing slide with optional logo."""
    slide = _blank(prs)
    _bg(slide, theme["bg"])

    _rect(slide, Inches(0), Inches(0),    Inches(0.18), SLIDE_H, theme["accent"])
    _rect(slide, Inches(0), Inches(3.45), SLIDE_W,      Inches(0.07), theme["accent2"])

    # Logo centred top
    if logo_bytes:
        try:
            img    = PILImage.open(io.BytesIO(logo_bytes))
            if img.mode not in ("RGB", "RGBA"):
                img = img.convert("RGBA")
            iw, ih = img.size
            max_w  = Inches(1.4)
            max_h  = Inches(0.7)
            scale  = min(max_w / iw, max_h / ih)
            fw     = int(iw * scale)
            fh     = int(ih * scale)
            buf    = io.BytesIO()
            img.save(buf, format="PNG")
            buf.seek(0)
            left_pos = int((SLIDE_W - fw) / 2)
            slide.shapes.add_picture(buf, left=left_pos, top=Inches(0.45), width=fw, height=fh)
        except Exception as e:
            print(f"[pptx] Closing logo skipped: {e}")

    _textbox(slide,
             left=Inches(1.5), top=Inches(1.5),
             width=Inches(10.5), height=Inches(1.8),
             text="Thank You", size=54, bold=True,
             color=theme["accent"], align=PP_ALIGN.CENTER)

    _textbox(slide,
             left=Inches(1.5), top=Inches(3.6),
             width=Inches(10.5), height=Inches(1.0),
             text=f"Let's build the future of {startup_name} together.",
             size=22, color=theme["text"], align=PP_ALIGN.CENTER)

    _textbox(slide,
             left=Inches(1.5), top=Inches(5.0),
             width=Inches(10.5), height=Inches(0.5),
             text="CONFIDENTIAL  ·  INVESTOR PRESENTATION",
             size=10, color=theme["muted"], align=PP_ALIGN.CENTER)


# ── Public entry point ────────────────────────────────────────────────────────

def generate_ppt_file(slides_data: list, context: dict = None) -> io.BytesIO:
    """
    Generate a PowerPoint file from slide data.
    Drop-in replacement — same signature as original.

    Args:
        slides_data : List of dicts with 'title', 'content', 'chart_data' keys
        context     : Optional startup context dict with keys:
                        theme       → "dark" | "light" | "blue"  (default: "dark")
                        startupName → used for footer watermark + logo fetch
                        logoUrl     → optional direct URL to company logo image

    Returns:
        BytesIO object containing the PPTX file, seeked to position 0
    """
    context      = context or {}
    theme        = _get_theme(context)
    startup_name = context.get("startupName", "")

    # ── Auto logo fetch ───────────────────────────────────────────────────────
    # Priority: explicit logoUrl in context > multi-source auto-fetch by name
    logo_bytes = None
    logo_url   = context.get("logoUrl", "").strip()
    if logo_url:
        try:
            r = requests.get(logo_url, timeout=5)
            if r.status_code == 200 and "image" in r.headers.get("content-type", ""):
                logo_bytes = r.content
                print(f"[pptx] ✓ Logo loaded from logoUrl")
        except Exception as e:
            print(f"[pptx] logoUrl fetch failed: {e}")

    if not logo_bytes and startup_name:
        logo_bytes = _fetch_logo(startup_name)

    # ── Build presentation ────────────────────────────────────────────────────
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H

    for i, slide_data in enumerate(slides_data):
        if not isinstance(slide_data, dict):
            continue
        if i == 0:
            _build_title_slide(prs, slide_data, theme, startup_name, logo_bytes)
        else:
            _build_content_slide(prs, slide_data,
                                 slide_num=i + 1,
                                 theme=theme,
                                 startup_name=startup_name)

    # Always append a Thank You closing slide
    _build_closing_slide(prs, theme, startup_name, logo_bytes)

    output = io.BytesIO()
    prs.save(output)
    output.seek(0)
    return output
